##from pptx import Presentation
##prs = Presentation("MAImatch.pptx")
##for slide in prs.slides:
##    for shape in slide.shapes:
##        if not shape.has_text_frame:
##            continue
##        text_frame = shape.text_frame
##        print(text_frame.paragraphs[0])

from pptx import Presentation
import pandas as pd
import nltk
from nltk.tokenize import sent_tokenize
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import re, string
from nltk.tag import pos_tag
from nltk.stem.wordnet import WordNetLemmatizer

def remove_noise(tweet_tokens, stop_words = ()):

    cleaned_tokens = []

    for token, tag in pos_tag(tweet_tokens):
        token = re.sub('http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+#]|[!*\(\),]|'\
                       '(?:%[0-9a-fA-F][0-9a-fA-F]))+','', token)
        token = re.sub("(@[A-Za-z0-9_]+)","", token)

        if tag.startswith("NN"):
            pos = 'n'
        elif tag.startswith('VB'):
            pos = 'v'
        else:
            pos = 'a'

        lemmatizer = WordNetLemmatizer()
        token = lemmatizer.lemmatize(token, pos)

        if len(token) > 0 and token not in string.punctuation and token.lower() not in stop_words:
            cleaned_tokens.append(token.lower())
    return cleaned_tokens


df = pd.read_excel('BOMReport.xlsx')
print("Column headings:")
#print(df)
IPN = df["Unnamed: 4"]
IPNN = IPN.tolist()


cleanedIPNN = [x for x in IPNN if str(x) != 'nan']

prs = Presentation('MAImatch.pptx')
##fo = open("content_steve.txt", "w")
# text_runs will be populated with a list of strings,
# one for each text run in presentation
list_of_elements = []
for slide in prs.slides:

	for shape in slide.shapes:
		if not shape.has_text_frame:
			continue
		for paragraph in shape.text_frame.paragraphs:
			line = ''
			for run in paragraph.runs:
				line += run.text
		list_of_elements.append(line)

#print(list_of_elements)
list_of_clean_word = []
for x in list_of_elements:
    if x != '':
        list_of_clean_word.append(x)
#print(list_of_clean_word)
cleanword = []
for x in list_of_clean_word:
    try:
        convt = int(x)
    except Exception:
        cleanword.append(x)
#print(cleanword)

filteredbasedonIPN = []
for x in cleanword:
    tokenized_word=word_tokenize(x)
    #print(tokenized_word)
    flag = 1
    for y in tokenized_word:
        if y in cleanedIPNN:
            if flag == 1:
                flag = 0
                filteredbasedonIPN.append(x)
    #print(filteredbasedonIPN)
#print(filteredbasedonIPN)

stop_words=set(stopwords.words("english"))
filtered_sent=[]
for w in filteredbasedonIPN:
    if w not in stop_words:
        filtered_sent.append(w)
##    print("Tokenized Sentence:",filteredbasedonIPN)
##    print("Filterd Sentence:",w)
print(filtered_sent)
def concatenate_list_data(list):
    result= ''
    for element in list:
        result += str(element)
    return result


k = concatenate_list_data(filtered_sent)
print(k)
#print(remove_noise(filtered_sent, stop_words))
#doc =
#print(filtered_sent)
##Consumed = []
##IPNNN = []
##for x in filtered_sent:
##    tokens=nltk.word_tokenize(x)
##    y = nltk.pos_tag(tokens)
####    print(x)
##    #print(y)
##    whatwewabr = []
##    howmany = []
##    for i in y:
##        name = i[0]
##        if name in cleanedIPNN:
##            whatwewabr.append(name)
##        typeof = i[1]
##        if typeof == "CD":
##            howmany.append(name)
##    for i in range (len(howmany)):
##        split = list(howmany[i])
##        if len(split) == 2:
##            for j in split:
##                try:
##                    cont = int(j)
##                    Consumed.append(cont)
##                    IPNNN.append(whatwewabr[i])
##                except Exception:
##                    pass
##        else:
##            Consumed.append(1)
##            IPNNN.append(whatwewabr[i])
##print(Consumed)
##print(IPNNN)


##    for i in cleanedIPNN:





##    tokenized_text=sent_tokenize(x)
##    print(tokenized_text)
##for i in list_of_elements:
##    j = type(i)
##    if j  == str:
##        print(i)

    #print(type(i))
##	for elements in list_of_elements:
##		fo.write(elements +'\n')
##	fo.write('\n')
##fo.close()

